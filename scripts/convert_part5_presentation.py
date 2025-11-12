#!/usr/bin/env python3
"""
Convert part5-vendor-lock-in.md to PowerPoint format.
Handles Mermaid diagrams, tables, and enhanced formatting.
"""

import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Constants
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
BLUE = RGBColor(65, 105, 225)  # Royal Blue for titles
GREEN = RGBColor(34, 139, 34)  # Forest Green for positives
ORANGE = RGBColor(255, 140, 0)  # Dark Orange for warnings
RED = RGBColor(220, 20, 60)  # Crimson for critical
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(80, 80, 80)
LIGHT_GRAY = RGBColor(150, 150, 150)

def create_prs():
    """Create a new presentation with custom slide size."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs

def add_title_slide(prs, title, subtitle=None):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(11.333), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = BLUE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4), Inches(11.333), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = GRAY
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

def parse_table(table_text):
    """Parse markdown table and return formatted text."""
    lines = [line.strip() for line in table_text.split('\n') if line.strip() and not line.strip().startswith('|---')]

    result = []
    for i, line in enumerate(lines):
        if line.startswith('|'):
            # Remove leading/trailing pipes and split
            cells = [cell.strip() for cell in line.strip('|').split('|')]

            # Skip separator lines
            if all(set(cell.strip()) <= {'-', ' ', ':'} for cell in cells):
                continue

            # First row is header
            if i == 0:
                result.append(' | '.join(cells))
                result.append('─' * 60)
            else:
                # Format regular rows
                # Remove HTML breaks and clean up
                cells = [cell.replace('<br/>', ' / ') for cell in cells]
                result.append(' | '.join(cells))

    return '\n'.join(result)

def extract_ascii_from_details(section_text):
    """Extract ASCII diagram from <details> block."""
    # Find content between <details> and </details>
    details_match = re.search(r'<details>.*?<summary>.*?</summary>\s*```text\s*(.+?)\s*```\s*</details>',
                              section_text, re.DOTALL)
    if details_match:
        return details_match.group(1).strip()
    return None

def format_content_line(p, line, level=0):
    """Format a single line of content."""
    line = line.strip()

    # Handle emojis and status indicators
    if line.startswith('✅') or line.startswith('✓') or '✅' in line[:5]:
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = GREEN
        p.level = level
    elif line.startswith('⚠️') or '⚠️' in line[:5] or line.startswith('🟡'):
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = ORANGE
        p.level = level
    elif line.startswith('🔴') or '🔴' in line[:5] or line.startswith('❌'):
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = RED
        p.level = level
    elif line.startswith('•') or line.startswith('-') or line.startswith('├─') or line.startswith('└─'):
        # Bullet point
        clean_line = line.replace('├─', '•').replace('└─', '•').replace('  ', ' ')
        if clean_line.startswith('-'):
            clean_line = '•' + clean_line[1:]
        p.text = clean_line
        p.font.size = Pt(14)
        p.font.color.rgb = BLACK
        p.level = level + 1
    elif line.startswith('│'):
        # Tree structure
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY
        p.level = level
    elif line.startswith('>'):
        # Blockquote
        p.text = line.replace('>', '').strip()
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = BLUE
        p.level = level
    else:
        # Clean markdown formatting
        clean_line = re.sub(r'\*\*(.+?)\*\*', r'\1', line)
        clean_line = re.sub(r'\[(.+?)\]\(.+?\)', r'\1', clean_line)
        p.text = clean_line

        # Bold headers
        if line.startswith('**') or (line and line[0].isupper() and line.endswith(':')):
            p.font.size = Pt(18)
            p.font.bold = True
        else:
            p.font.size = Pt(14)
        p.font.color.rgb = BLACK
        p.level = level

def add_content_slide(prs, title, content_text, speaker_notes=None):
    """Add a content slide with title and body text."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7)
    )
    title_frame = title_box.text_frame
    # Clean title
    clean_title = title.replace('### Slide ', '').replace('Slide ', '')
    clean_title = re.sub(r':\s*.*$', '', clean_title)  # Remove subtitle after colon
    title_frame.text = clean_title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = BLUE

    # Content
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.1), Inches(12.333), Inches(5.9)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    # Parse and format content
    lines = content_text.split('\n')
    for i, line in enumerate(lines):
        if not line.strip():
            continue

        if i > 0:
            content_frame.add_paragraph()

        p = content_frame.paragraphs[-1]
        format_content_line(p, line)

    # Add speaker notes
    if speaker_notes:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        # Clean up speaker notes
        clean_notes = re.sub(r'<details>.*?</details>', '', speaker_notes, flags=re.DOTALL)
        clean_notes = re.sub(r'\*\*\[Reference:.*?\]\*\*', '', clean_notes)
        clean_notes = re.sub(r'\n\n+', '\n\n', clean_notes)
        text_frame.text = clean_notes.strip()

    return slide

def parse_slide_section(section_text):
    """Parse a slide section and extract title, content, and notes."""
    lines = section_text.split('\n')

    # Extract title
    title = None
    for line in lines:
        if line.startswith('### Slide'):
            title = line.strip()
            break

    if not title:
        return None, None, None

    # Try to extract ASCII diagram from details block
    ascii_diagram = extract_ascii_from_details(section_text)

    # Extract tables
    content_parts = []

    # Look for tables
    table_matches = re.finditer(r'(\|.+?\|(?:\n\|.+?\|)+)', section_text, re.MULTILINE)
    for match in table_matches:
        table_text = match.group(1)
        formatted_table = parse_table(table_text)
        if formatted_table:
            content_parts.append(formatted_table)

    # If we have ASCII diagram, use it
    if ascii_diagram:
        content_parts.insert(0, ascii_diagram)

    # Look for assessment summaries, options, etc.
    # Extract blockquotes
    blockquotes = re.findall(r'^>\s*(.+?)$', section_text, re.MULTILINE)
    if blockquotes:
        content_parts.append('\n'.join(['> ' + bq for bq in blockquotes]))

    # Extract bullet lists that aren't in tables
    in_table = False
    list_items = []
    for line in lines:
        if line.strip().startswith('|'):
            in_table = True
            continue
        if in_table and not line.strip().startswith('|'):
            in_table = False

        if not in_table:
            if line.strip().startswith(('- ', '• ', '├─', '└─')) or re.match(r'^\s*[🔴⚠️✅🟡]', line):
                list_items.append(line)

    if list_items and not content_parts:
        content_parts.append('\n'.join(list_items))

    # Join all content
    content_text = '\n\n'.join(content_parts)

    # Extract speaker notes
    speaker_notes = None
    if '**Speaker Notes:**' in section_text:
        notes_section = section_text.split('**Speaker Notes:**')[1]
        notes_section = notes_section.split('---')[0]
        notes_section = notes_section.split('###')[0]
        speaker_notes = notes_section.strip()

    return title, content_text, speaker_notes

def main():
    """Main function to convert part5-vendor-lock-in.md to PowerPoint."""
    base_path = Path('/Users/bastiengallay/Dev/clients/CISAC/docs/deliverables/first-restitution-2025-11-24')
    md_path = base_path / 'presentation' / 'part5-vendor-lock-in.md'
    output_path = base_path / 'powerpoint' / 'part5-vendor-lock-in.pptx'

    # Create output directory if it doesn't exist
    output_path.parent.mkdir(parents=True, exist_ok=True)

    print("=" * 60)
    print("Part 5: Vendor Lock-in PowerPoint Converter")
    print("=" * 60)
    print()
    print(f"Input:  {md_path.name}")
    print(f"Output: {output_path}")
    print()

    # Read markdown
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Create presentation
    prs = create_prs()

    # Add title slide
    part_title = "PART 5: Vendor Lock-in Reality"
    add_title_slide(prs, part_title, "7 minutes, Slides 18-21")

    # Split into slide sections
    slide_sections = re.split(r'(?=### Slide \d+:)', content)

    slide_count = 0
    for section in slide_sections:
        if not section.strip() or '### Slide' not in section:
            continue

        title, content_text, speaker_notes = parse_slide_section(section)

        if title and content_text:
            add_content_slide(prs, title, content_text, speaker_notes)
            slide_count += 1
            print(f"  ✓ Added: {title}")

    # Save
    prs.save(output_path)

    print()
    print("=" * 60)
    print("COMPLETE")
    print("=" * 60)
    size_mb = output_path.stat().st_size / (1024 * 1024)
    print(f"Created: {output_path.name}")
    print(f"Slides:  {len(prs.slides)} total ({slide_count} content slides)")
    print(f"Size:    {size_mb:.2f} MB")
    print()

if __name__ == '__main__':
    main()
