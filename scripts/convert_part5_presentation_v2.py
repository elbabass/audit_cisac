#!/usr/bin/env python3
"""
Convert part5-vendor-lock-in.md to PowerPoint format.
Enhanced version that properly handles tables, diagrams, and new formatting.
"""

import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Constants
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
BLUE = RGBColor(65, 105, 225)
GREEN = RGBColor(34, 139, 34)
ORANGE = RGBColor(255, 140, 0)
RED = RGBColor(220, 20, 60)
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(80, 80, 80)

def create_prs():
    """Create a new presentation with custom slide size."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs

def add_title_slide(prs, title, subtitle=None):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(11.333), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = BLUE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4), Inches(11.333), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = GRAY
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

def format_text_run(run, text):
    """Apply formatting to a text run based on content."""
    # Clean markdown
    clean_text = text.replace('**', '')
    run.text = clean_text

    # Detect bold
    if '**' in text:
        run.font.bold = True

def add_paragraph_with_formatting(text_frame, text, level=0, font_size=14, color=BLACK, bold=False):
    """Add a formatted paragraph."""
    if len(text_frame.paragraphs) == 1 and not text_frame.paragraphs[0].text:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()

    # Clean text
    clean_text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    clean_text = clean_text.replace('<br/>', ' ')

    p.text = clean_text
    p.level = level
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold

    return p

def extract_slide_content(section_text):
    """Extract structured content from a slide section."""
    content = {
        'ascii_diagram': None,
        'tables': [],
        'bullets': [],
        'quotes': [],
        'text_blocks': []
    }

    # Extract ASCII from details block
    details_match = re.search(
        r'<details>.*?```text\s*(.+?)\s*```.*?</details>',
        section_text, re.DOTALL
    )
    if details_match:
        content['ascii_diagram'] = details_match.group(1).strip()

    # Extract tables (between mermaid blocks and speaker notes)
    # Find content between mermaid block and speaker notes
    main_content = section_text
    if '**Speaker Notes:**' in main_content:
        main_content = main_content.split('**Speaker Notes:**')[0]

    # Skip mermaid blocks
    main_content = re.sub(r'```mermaid.+?```', '', main_content, flags=re.DOTALL)

    # Extract tables
    table_pattern = r'\|[^\n]+\|(?:\n\|[^\n]+\|)+'
    for match in re.finditer(table_pattern, main_content):
        table_text = match.group(0)
        # Parse table
        rows = []
        for line in table_text.split('\n'):
            if '|' in line and not re.match(r'^\|[\s\-:]+\|', line):
                cells = [c.strip() for c in line.strip('|').split('|')]
                rows.append(cells)
        if rows:
            content['tables'].append(rows)

    # Extract blockquotes
    for match in re.finditer(r'^>\s*(.+?)$', main_content, re.MULTILINE):
        content['quotes'].append(match.group(1))

    # Extract bullet points (outside of details and mermaid)
    bullet_section = main_content
    bullet_section = re.sub(r'<details>.*?</details>', '', bullet_section, flags=re.DOTALL)

    for line in bullet_section.split('\n'):
        line = line.strip()
        if line.startswith(('- ', '• ', '├─', '└─')) or re.match(r'^[🔴⚠️✅🟡]', line):
            content['bullets'].append(line)

    return content

def add_content_slide(prs, title, section_text):
    """Add a content slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7)
    )
    title_frame = title_box.text_frame
    clean_title = title.replace('### Slide ', '').replace('Slide ', '')
    clean_title = re.sub(r'\s*-\s*.+$', '', clean_title)  # Remove subtitle
    title_frame.text = clean_title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = BLUE

    # Extract content
    content = extract_slide_content(section_text)

    # Content area
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.1), Inches(12.333), Inches(5.9)
    )
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    # Add ASCII diagram first if available
    if content['ascii_diagram']:
        lines = content['ascii_diagram'].split('\n')[:25]  # Limit lines
        for line in lines:
            add_paragraph_with_formatting(text_frame, line, level=0, font_size=11, color=GRAY)

    # Add tables
    elif content['tables']:
        for table in content['tables']:
            for i, row in enumerate(table[:15]):  # Limit rows
                row_text = ' | '.join(row)
                # Clean HTML
                row_text = row_text.replace('<br/>', ' / ')

                if i == 0:
                    # Header row
                    add_paragraph_with_formatting(text_frame, row_text, level=0, font_size=14, color=BLUE, bold=True)
                else:
                    # Determine color based on content
                    color = BLACK
                    if '✅' in row_text or 'GRANTED' in row_text or 'POSITIVE' in row_text:
                        color = GREEN
                    elif '🔴' in row_text or 'BLOCKED' in row_text or 'EXCLUDED' in row_text:
                        color = RED
                    elif '⚠️' in row_text or '🟡' in row_text or 'PENDING' in row_text or 'LIMITED' in row_text:
                        color = ORANGE

                    add_paragraph_with_formatting(text_frame, row_text, level=0, font_size=12, color=color)

    # Add quotes
    elif content['quotes']:
        for quote in content['quotes']:
            add_paragraph_with_formatting(text_frame, quote, level=0, font_size=16, color=BLUE, bold=True)

    # Add bullets
    elif content['bullets']:
        for bullet in content['bullets'][:20]:  # Limit bullets
            # Determine color and size
            color = BLACK
            size = 14

            if bullet.startswith('✅') or '✅' in bullet[:5]:
                color = GREEN
                size = 15
            elif bullet.startswith('🔴') or '🔴' in bullet[:5]:
                color = RED
                size = 15
            elif bullet.startswith('⚠️') or '⚠️' in bullet[:5] or bullet.startswith('🟡'):
                color = ORANGE
                size = 15

            # Clean bullet
            clean_bullet = bullet.replace('├─', '•').replace('└─', '•')
            if clean_bullet.startswith('- '):
                clean_bullet = '• ' + clean_bullet[2:]

            level = 1 if bullet.startswith(('  ', '\t')) else 0
            add_paragraph_with_formatting(text_frame, clean_bullet, level=level, font_size=size, color=color)

    # Add speaker notes
    if '**Speaker Notes:**' in section_text:
        notes_text = section_text.split('**Speaker Notes:**')[1]
        notes_text = notes_text.split('---')[0].split('###')[0]

        # Include ASCII diagram in notes
        if content['ascii_diagram']:
            notes_text = "DIAGRAM:\n" + content['ascii_diagram'] + "\n\n" + notes_text

        # Clean notes
        notes_text = re.sub(r'<details>.*?</details>', '', notes_text, flags=re.DOTALL)
        notes_text = re.sub(r'\*\*\[Reference:.*?\]\*\*', '', notes_text)
        notes_text = re.sub(r'\n\n+', '\n\n', notes_text).strip()

        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text

    return slide

def main():
    """Main conversion function."""
    base_path = Path('/Users/bastiengallay/Dev/clients/CISAC/docs/deliverables/first-restitution-2025-11-24')
    md_path = base_path / 'presentation' / 'part5-vendor-lock-in.md'
    output_path = base_path / 'powerpoint' / 'part5-vendor-lock-in.pptx'

    output_path.parent.mkdir(parents=True, exist_ok=True)

    print("=" * 60)
    print("Part 5: Vendor Lock-in PowerPoint Converter (v2)")
    print("=" * 60)
    print(f"\nInput:  {md_path}")
    print(f"Output: {output_path}\n")

    # Read markdown
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Create presentation
    prs = create_prs()
    add_title_slide(prs, "PART 5: Vendor Lock-in Reality", "7 minutes, Slides 18-21")

    # Process slides
    slides = re.split(r'(?=### Slide \d+:)', content)

    for section in slides:
        if '### Slide' not in section:
            continue

        # Extract title
        title_match = re.search(r'### (Slide \d+:.+?)$', section, re.MULTILINE)
        if title_match:
            title = title_match.group(1)
            add_content_slide(prs, title, section)
            print(f"  ✓ {title}")

    # Save
    prs.save(output_path)

    size_mb = output_path.stat().st_size / (1024 * 1024)
    print(f"\n{'=' * 60}")
    print("COMPLETE")
    print('=' * 60)
    print(f"File:   {output_path.name}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Size:   {size_mb:.2f} MB\n")

if __name__ == '__main__':
    main()
