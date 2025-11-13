#!/usr/bin/env python3
"""
Generate PowerPoint presentation for Part 1 (Journey) from markdown source.
Applies Teragone-inspired design with CISAC branding.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
import os

# Design System Colors (Teragone brand colors)
TERAGONE_DARK_BLUE = RGBColor(6, 38, 67)  # #062643 - Bleu Marine Teragone (hero background)
TERAGONE_BLUE = RGBColor(0, 47, 108)  # #002f6c - Primary blue
TERAGONE_LIGHT_BLUE = RGBColor(41, 98, 155)  # #29629b - Lighter blue for accents
CISAC_RED = RGBColor(226, 0, 38)  # #e20026 - Only for CISAC logo
WHITE = RGBColor(255, 255, 255)  # #ffffff
LIGHT_GRAY = RGBColor(245, 245, 245)  # #f5f5f5
DARK_GRAY = RGBColor(64, 64, 64)  # #404040
MEDIUM_GRAY = RGBColor(128, 128, 128)  # #808080

# Paths
BASE_DIR = "/Users/bastiengallay/Dev/clients/CISAC"
CISAC_LOGO_PATH = f"{BASE_DIR}/docs/resources/CISAC-logo.png"
TERAGONE_LOGO_PATH = f"{BASE_DIR}/docs/resources/Blanc_Teragone-Factory-logo.png"
DIAGRAM_DIR = f"{BASE_DIR}/docs/deliverables/first-restitution-2025-11-24/presentation/diagrams"
OUTPUT_PATH = f"{BASE_DIR}/docs/deliverables/first-restitution-2025-11-24/powerpoint/part1-journey-v2.pptx"


def create_cisac_logo_png():
    """Create a simple PNG representation of CISAC logo."""
    # Create a 200x200 red square with white CISAC text
    size = 200
    img = Image.new('RGB', (size, size), color=(226, 0, 38))  # CISAC red
    draw = ImageDraw.Draw(img)

    # Since we can't render the complex SVG paths easily without cairo,
    # we'll create a simple red square as placeholder
    # The actual logo will need to be exported from the SVG manually

    img.save(CISAC_LOGO_PATH)
    print(f"Created simple logo placeholder at {CISAC_LOGO_PATH}")


def add_title_slide(prs):
    """Create Slide 1: Title slide with dark blue hero background and dual logos."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background - Teragone dark blue hero
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = TERAGONE_DARK_BLUE

    # Add CISAC logo (top-left)
    if os.path.exists(CISAC_LOGO_PATH):
        cisac_logo = slide.shapes.add_picture(
            CISAC_LOGO_PATH,
            Inches(0.5), Inches(0.4),
            height=Inches(0.7)
        )

    # Add Teragone logo (top-right)
    if os.path.exists(TERAGONE_LOGO_PATH):
        teragone_logo = slide.shapes.add_picture(
            TERAGONE_LOGO_PATH,
            Inches(7.8), Inches(0.4),
            height=Inches(0.7)
        )

    # Main title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(8), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = "ISWC System Audit"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE  # White text on dark background
    title_para.alignment = PP_ALIGN.LEFT

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.6),
        Inches(8), Inches(0.6)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "First Conclusions & Strategic Recommendations"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = LIGHT_GRAY  # Light gray for subtitle on dark
    subtitle_para.alignment = PP_ALIGN.LEFT

    # Team info
    team_box = slide.shapes.add_textbox(
        Inches(1), Inches(5),
        Inches(8), Inches(1)
    )
    team_frame = team_box.text_frame
    team_frame.text = "Teragone-Factory\nGuillaume Jay & Bastien Gallay"
    for para in team_frame.paragraphs:
        para.font.size = Pt(20)
        para.font.color.rgb = WHITE  # White text on dark background
        para.alignment = PP_ALIGN.LEFT

    # Date
    date_box = slide.shapes.add_textbox(
        Inches(1), Inches(6.2),
        Inches(8), Inches(0.5)
    )
    date_frame = date_box.text_frame
    date_frame.text = "November 24, 2025"
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(18)
    date_para.font.color.rgb = LIGHT_GRAY  # Light gray on dark background
    date_para.alignment = PP_ALIGN.LEFT


def add_header_branding(slide, slide_number, total_slides, part_name="Part 1"):
    """Add consistent branding header to content slides."""
    # Thin accent bar at top
    accent_bar = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(0.05)
    )
    accent_fill = accent_bar.fill
    accent_fill.solid()
    accent_fill.fore_color.rgb = TERAGONE_DARK_BLUE

    # Slide number in header (top-right)
    number_box = slide.shapes.add_textbox(
        Inches(8), Inches(0.1),
        Inches(1.5), Inches(0.3)
    )
    number_frame = number_box.text_frame
    number_frame.text = f"{part_name} | {slide_number}/{total_slides}"
    number_para = number_frame.paragraphs[0]
    number_para.font.size = Pt(10)
    number_para.font.color.rgb = MEDIUM_GRAY
    number_para.alignment = PP_ALIGN.RIGHT


def add_journey_slide(prs):
    """Create Slide 2: The Audit Journey timeline."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Add header branding
    add_header_branding(slide, 2, 3, "Part 1")

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "The Audit Journey - How We Got Here"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = TERAGONE_BLUE

    # Add diagram if exists
    diagram_path = f"{DIAGRAM_DIR}/part1-journey_1_.png"
    if os.path.exists(diagram_path):
        slide.shapes.add_picture(
            diagram_path,
            Inches(0.5), Inches(1.1),
            width=Inches(5.5)
        )

    # Key Statistics table (right side)
    stats_box = slide.shapes.add_textbox(
        Inches(6.2), Inches(1.1),
        Inches(3.5), Inches(3)
    )
    stats_frame = stats_box.text_frame
    stats_frame.word_wrap = True

    # Title for stats
    p = stats_frame.paragraphs[0]
    p.text = "Key Statistics"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TERAGONE_LIGHT_BLUE

    # Stats content
    stats = [
        ("Budget:", "20 man-days"),
        ("Workshops:", "6+ sessions"),
        ("Azure Resources:", "343 analyzed"),
        ("Code Review:", "42+ files"),
        ("Documentation:", "95+ validation rules")
    ]

    for label, value in stats:
        p = stats_frame.add_paragraph()
        p.text = f"{label} {value}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)

    # Footer note
    note_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.8),
        Inches(9), Inches(0.5)
    )
    note_frame = note_box.text_frame
    note_frame.text = "From technical assessment to strategic vendor independence evaluation"
    note_para = note_frame.paragraphs[0]
    note_para.font.size = Pt(12)
    note_para.font.italic = True
    note_para.font.color.rgb = DARK_GRAY


def add_pivot_slide(prs):
    """Create Slide 3: The Strategic Pivot."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Add header branding
    add_header_branding(slide, 3, 3, "Part 1")

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "The Strategic Pivot - What This Audit Became"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = TERAGONE_BLUE

    # Two-column comparison
    # Left column
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2),
        Inches(4.5), Inches(4)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    p = left_frame.paragraphs[0]
    p.text = "What We Were Asked To Do"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TERAGONE_LIGHT_BLUE

    left_items = [
        "✓ Validate Hyperscale proposal",
        "✓ Assess architecture quality",
        "✓ Identify technical debt",
        "✓ Migration roadmap"
    ]

    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)

    # Right column
    right_box = slide.shapes.add_textbox(
        Inches(5.5), Inches(1.2),
        Inches(4.5), Inches(4)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    p = right_frame.paragraphs[0]
    p.text = "What We Actually Discovered"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TERAGONE_LIGHT_BLUE

    right_items = [
        ("⚠️", "Cost control gap (€600K/year)"),
        ("✓", "Platform technically sound"),
        ("✓", "Recently upgraded, maintained"),
        ("⚠️", "Governance gaps - real issue"),
        ("🔴", "Vendor independence unknown"),
        ("🔴", "Knowledge transfer HIGH RISK")
    ]

    for emoji, text in right_items:
        p = right_frame.add_paragraph()
        p.text = f"{emoji} {text}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)

    # Key insight box at bottom
    insight_box = slide.shapes.add_textbox(
        Inches(2), Inches(5.5),
        Inches(6), Inches(1.2)
    )
    # Add border-like background
    shape_fill = insight_box.fill
    shape_fill.solid()
    shape_fill.fore_color.rgb = LIGHT_GRAY

    insight_frame = insight_box.text_frame
    insight_frame.word_wrap = True
    insight_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = insight_frame.paragraphs[0]
    p.text = '"The platform is solid. Control is the problem."'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.italic = True
    p.font.color.rgb = TERAGONE_BLUE  # Key insight in Teragone blue
    p.alignment = PP_ALIGN.CENTER


def main():
    """Generate the PowerPoint presentation."""
    print("Generating Part 1 (Journey) PowerPoint presentation...")

    # Create logo if it doesn't exist
    if not os.path.exists(CISAC_LOGO_PATH):
        create_cisac_logo_png()

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add slides
    print("Creating Slide 1: Title slide...")
    add_title_slide(prs)

    print("Creating Slide 2: Journey timeline...")
    add_journey_slide(prs)

    print("Creating Slide 3: Strategic pivot...")
    add_pivot_slide(prs)

    # Save presentation
    prs.save(OUTPUT_PATH)
    print(f"\n✓ Presentation saved to: {OUTPUT_PATH}")
    print(f"  - 3 slides generated")
    print(f"  - Teragone design applied")
    print(f"  - CISAC branding included")


if __name__ == "__main__":
    main()
