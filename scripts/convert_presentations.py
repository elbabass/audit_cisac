#!/usr/bin/env python3
"""
Convert markdown presentation files to PowerPoint (.pptx) format.
Handles presentation parts 2-7 and annexes A-D.
"""

import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Constants
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
BLUE = RGBColor(65, 105, 225)  # Royal Blue for titles
GREEN = RGBColor(34, 139, 34)  # Forest Green for positives
ORANGE = RGBColor(255, 140, 0)  # Dark Orange for warnings
RED = RGBColor(220, 20, 60)  # Crimson for critical
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(80, 80, 80)

def create_prs():
    """Create a new presentation with custom slide size."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs

def add_title_slide(prs, title, subtitle=None):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(11.333), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = BLUE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4), Inches(11.333), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = GRAY
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_text, speaker_notes=None):
    """Add a content slide with title and body text."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = BLUE

    # Content
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.1), Inches(12.333), Inches(6)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    # Parse and format content
    lines = content_text.split('\n')
    for i, line in enumerate(lines):
        if i > 0:
            content_frame.add_paragraph()

        p = content_frame.paragraphs[-1]

        # Handle different line types
        if line.strip().startswith('✅') or line.strip().startswith('✓'):
            p.text = line.strip()
            p.font.size = Pt(16)
            p.font.color.rgb = GREEN
            p.level = 0
        elif line.strip().startswith('⚠️') or line.strip().startswith('🟠'):
            p.text = line.strip()
            p.font.size = Pt(16)
            p.font.color.rgb = ORANGE
            p.level = 0
        elif line.strip().startswith('🔴') or line.strip().startswith('❌'):
            p.text = line.strip()
            p.font.size = Pt(16)
            p.font.color.rgb = RED
            p.level = 0
        elif line.strip().startswith('•') or line.strip().startswith('├─') or line.strip().startswith('└─'):
            # Bullet point
            p.text = line.strip().replace('├─', '•').replace('└─', '•').replace('  ', '')
            p.font.size = Pt(14)
            p.font.color.rgb = BLACK
            p.level = 1 if line.startswith('  ') else 0
        else:
            p.text = line.strip()
            if line.strip() and line[0].isupper():
                p.font.size = Pt(18)
                p.font.bold = True
            else:
                p.font.size = Pt(14)
            p.font.color.rgb = BLACK

    # Add speaker notes
    if speaker_notes:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = speaker_notes

    return slide

def parse_markdown_slide(slide_text):
    """Parse a slide section from markdown."""
    lines = slide_text.split('\n')

    # Extract title (first ### line)
    title = None
    for line in lines:
        if line.startswith('### Slide '):
            title = line.replace('### ', '').strip()
            break

    if not title:
        return None, None, None

    # Extract content between **Visual:** or ``` blocks
    content = []
    in_code_block = False
    code_block_content = []

    for line in lines:
        if line.strip().startswith('```'):
            if in_code_block:
                # End of code block
                content.append('\n'.join(code_block_content))
                code_block_content = []
                in_code_block = False
            else:
                # Start of code block
                in_code_block = True
                # Skip language identifier line
                continue
        elif in_code_block:
            code_block_content.append(line)

    # Extract speaker notes
    speaker_notes = None
    if '**Speaker Notes:**' in slide_text:
        notes_section = slide_text.split('**Speaker Notes:**')[1]
        notes_section = notes_section.split('---')[0]  # Stop at next divider
        notes_section = notes_section.split('###')[0]  # Stop at next slide
        speaker_notes = notes_section.strip()
        # Clean up references
        speaker_notes = re.sub(r'\*\*\[Reference:.*?\]\*\*', '', speaker_notes)
        speaker_notes = re.sub(r'\n\n+', '\n\n', speaker_notes)

    content_text = '\n'.join(content) if content else ''

    return title, content_text, speaker_notes

def process_presentation_file(md_file_path):
    """Process a markdown presentation file and create PowerPoint."""
    md_path = Path(md_file_path)
    output_path = md_path.with_suffix('.pptx')

    print(f"Processing: {md_path.name}")

    # Read markdown content
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Create presentation
    prs = create_prs()

    # Extract part title from first # heading
    part_title_match = re.search(r'^# (.+?)$', content, re.MULTILINE)
    if part_title_match:
        part_title = part_title_match.group(1)
        # Clean up any markdown links
        part_title = re.sub(r'\[.*?\]\(.*?\)', '', part_title).strip()
        part_title = re.sub(r'\(.*? min.*?\)', '', part_title).strip()
        add_title_slide(prs, part_title)

    # Split content into slides
    slide_sections = re.split(r'\n### Slide \d+:', content)

    for section in slide_sections[1:]:  # Skip first (before any slide)
        section = '### Slide ' + section  # Add back the split marker

        title, content_text, speaker_notes = parse_markdown_slide(section)

        if title and content_text:
            add_content_slide(prs, title, content_text, speaker_notes)

    # Save presentation
    prs.save(output_path)
    print(f"  Created: {output_path.name} ({len(prs.slides)} slides)")

    return output_path

def process_annex_file(md_file_path):
    """Process an annex markdown file and create PowerPoint."""
    md_path = Path(md_file_path)
    output_path = md_path.with_suffix('.pptx')

    print(f"Processing annex: {md_path.name}")

    # Read markdown content
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Create presentation
    prs = create_prs()

    # Extract annex title
    annex_title_match = re.search(r'^# (.+?)$', content, re.MULTILINE)
    if annex_title_match:
        annex_title = annex_title_match.group(1)
        annex_title = re.sub(r'\[.*?\]\(.*?\)', '', annex_title).strip()
        add_title_slide(prs, annex_title)

    # For annexes, create slides based on ## sections
    sections = re.split(r'\n## ', content)

    for section in sections[1:]:  # Skip first (title)
        section = '## ' + section
        lines = section.split('\n')
        section_title = lines[0].replace('## ', '').strip()

        # Get content (skip first line which is title)
        section_content = '\n'.join(lines[1:])

        # Clean up markdown
        section_content = re.sub(r'\[.*?\]\(.*?\)', '', section_content)
        section_content = re.sub(r'\*\*(.+?)\*\*', r'\1', section_content)
        section_content = re.sub(r'\n\n+', '\n\n', section_content)

        # Limit content length for slides
        if len(section_content) > 1500:
            section_content = section_content[:1500] + '\n\n[Content truncated - see full documentation]'

        add_content_slide(prs, section_title, section_content)

    # Save presentation
    prs.save(output_path)
    print(f"  Created: {output_path.name} ({len(prs.slides)} slides)")

    return output_path

def main():
    """Main function to process all files."""
    base_path = Path('/Users/bastiengallay/Dev/clients/CISAC/docs/deliverables/first-restitution-2025-11-24')

    # Presentation files
    presentation_files = [
        base_path / 'presentation' / 'part2-executive-summary.md',
        base_path / 'presentation' / 'part3-technical-findings.md',
        base_path / 'presentation' / 'part4-governance-findings.md',
        base_path / 'presentation' / 'part5-vendor-lock-in.md',
        base_path / 'presentation' / 'part6-strategic-recommendations.md',
        base_path / 'presentation' / 'part7-decision-points.md',
    ]

    # Annex files
    annex_files = [
        base_path / 'annexes' / 'annex-a-detailed-findings.md',
        base_path / 'annexes' / 'annex-b-stakeholder-quotes.md',
        base_path / 'annexes' / 'annex-c-validation-rules.md',
        base_path / 'annexes' / 'annex-d-items-verification.md',
    ]

    print("=" * 60)
    print("ISWC Audit Presentation Converter")
    print("=" * 60)
    print()

    # Process presentation files
    print("Processing presentation parts...")
    print("-" * 60)
    output_files = []
    for md_file in presentation_files:
        if md_file.exists():
            output_file = process_presentation_file(md_file)
            output_files.append(output_file)
        else:
            print(f"  WARNING: File not found: {md_file.name}")

    print()
    print("Processing annexes...")
    print("-" * 60)
    for md_file in annex_files:
        if md_file.exists():
            output_file = process_annex_file(md_file)
            output_files.append(output_file)
        else:
            print(f"  WARNING: File not found: {md_file.name}")

    print()
    print("=" * 60)
    print("SUMMARY")
    print("=" * 60)
    print(f"Total files created: {len(output_files)}")
    print()
    print("Output files:")
    for output_file in output_files:
        size_mb = output_file.stat().st_size / (1024 * 1024)
        print(f"  - {output_file.name} ({size_mb:.2f} MB)")
    print()
    print("Done!")

if __name__ == '__main__':
    main()
